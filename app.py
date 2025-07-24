import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
import openai
import time

# Premium features configuration
PREMIUM_FEATURES = {
    "investor_contact_info": True,
    "detailed_funding_reports": True,
    "competitor_analysis": True,
    "market_size_estimates": True,
    "export_reports": True
}

# Configure Streamlit page with clean light mode
st.set_page_config(
    page_title="cofounder.ai | Startup Intelligence Platform",
    page_icon="🚀",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Simplified CSS styling for light mode
st.markdown("""
<style>
    /* Simplified light mode styling */
    .stApp {
        background-color: #ffffff;
        color: #333333;
    }
    
    /* Clean header */
    .header {
        background: linear-gradient(135deg, #4b6cb7 0%, #182848 100%);
        color: white;
        padding: 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        box-shadow: 0 4px 20px rgba(0,0,0,0.1);
    }
    
    /* Content cards */
    .content-card {
        border-radius: 12px;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        padding: 1.5rem;
        margin-bottom: 1.5rem;
        background: white;
        border-left: 4px solid #4b6cb7;
        transition: transform 0.3s ease;
        color: #333333 !important;
    }
    
    .content-card:hover {
        transform: translateY(-3px);
    }
    
    /* Buttons */
    .stButton>button {
        background: linear-gradient(135deg, #4b6cb7 0%, #182848 100%);
        color: white !important;
        border: none;
        border-radius: 8px;
        padding: 0.5rem 1.5rem;
        font-weight: 600;
    }
    
    /* Tabs */
    .stTabs [data-baseweb="tab"] {
        padding: 0.75rem 1.5rem;
        font-weight: 500;
        color: #4b6cb7;
    }
    
    .stTabs [aria-selected="true"] {
        background: #4b6cb7;
        color: white !important;
    }
    
    /* Ensure all text is visible */
    body {
        color: #333333 !important;
    }
    
    h1, h2, h3, h4, h5, h6 {
        color: #182848 !important;
    }
    
    /* Subscription plans */
    .pricing-card {
        border-radius: 12px;
        padding: 1.5rem;
        margin: 1rem 0;
        background: white;
        box-shadow: 0 4px 12px rgba(0,0,0,0.08);
        border: 1px solid #e0e0e0;
        color: #333333 !important;
    }
    
    .featured-badge {
        background: #4b6cb7;
        color: white;
        padding: 0.25rem 1rem;
        border-radius: 20px;
        font-size: 0.8rem;
        font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'premium_user' not in st.session_state:
    st.session_state.premium_user = True  # Demo mode - all users get premium

# Configure OpenAI API
openai.api_key = st.secrets["OPENAI_API_KEY"]

# ========== CORE BUSINESS FUNCTIONS ==========
def generate_competitor_analysis(startup_idea, industry):
    """Generate detailed competitor analysis using AI"""
    with st.spinner("Analyzing competitors..."):
        try:
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a startup competitive intelligence analyst."},
                    {"role": "user", "content": f"""
                    Analyze the competitive landscape for this startup idea in the {industry} industry:
                    
                    {startup_idea}
                    
                    Provide:
                    1. Direct competitors with their funding status
                    2. Market positioning analysis
                    3. Competitive advantages
                    4. Potential threats
                    5. Market share estimates
                    """}
                ]
            )
            return response["choices"][0]["message"]["content"]
        except Exception as e:
            st.error(f"Error generating analysis: {e}")
            return None

def generate_funding_strategy(startup_idea, industry, stage):
    """Generate personalized funding strategy"""
    with st.spinner("Creating funding strategy..."):
        try:
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a startup funding strategist."},
                    {"role": "user", "content": f"""
                    Create a detailed funding strategy for this {industry} startup at {stage} stage:
                    
                    {startup_idea}
                    
                    Include:
                    1. Recommended funding sources
                    2. Ideal investor profile
                    3. Valuation benchmarks
                    4. Funding timeline
                    5. Key metrics to focus on
                    """}
                ]
            )
            return response["choices"][0]["message"]["content"]
        except Exception as e:
            st.error(f"Error generating strategy: {e}")
            return None

def generate_pitch_deck_review(deck_text):
    """Generate detailed pitch deck review"""
    with st.spinner("Analyzing pitch deck..."):
        try:
            response = openai.ChatCompletion.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "You are a pitch deck expert who has reviewed thousands of decks."},
                    {"role": "user", "content": f"""
                    Provide a detailed review of this pitch deck:
                    
                    {deck_text}
                    
                    Cover:
                    1. Strengths and weaknesses
                    2. Storytelling effectiveness
                    3. Financial projections quality
                    4. Design recommendations
                    5. Suggested improvements
                    """}
                ]
            )
            return response["choices"][0]["message"]["content"]
        except Exception as e:
            st.error(f"Error analyzing deck: {e}")
            return None

# ========== MAIN APP ==========
def main():
    """Main application interface"""
    
    # Clean header
    st.markdown("""
    <div class="header">
        <h1 style="color:white;margin:0;">cofounder.ai</h1>
        <p style="color:white;margin:0;font-size:1.2rem;">The AI-powered startup intelligence platform</p>
    </div>
    """, unsafe_allow_html=True)
    
    # Main tabs
    tab1, tab2, tab3 = st.tabs(["🏠 Dashboard", "💰 Fundraising", "📊 Analytics"])
    
    with tab1:
        st.header("Startup Intelligence Dashboard")
        
        # Quick analysis section
        with st.expander("🚀 Quick Startup Analysis", expanded=True):
            startup_idea = st.text_area("Describe your startup (2-3 sentences)", height=100)
            industry = st.selectbox("Industry", ["Tech", "FinTech", "HealthTech", "EdTech", "AI/ML", "E-commerce", "Other"])
            
            if st.button("Analyze My Startup"):
                if startup_idea.strip():
                    with st.spinner("Generating insights..."):
                        # Generate multiple analyses in parallel
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.markdown("""
                            <div class="content-card">
                                <h4>Competitive Landscape</h4>
                            """, unsafe_allow_html=True)
                            analysis = generate_competitor_analysis(startup_idea, industry)
                            if analysis:
                                st.write(analysis)
                            st.markdown("</div>", unsafe_allow_html=True)
                        
                        with col2:
                            st.markdown("""
                            <div class="content-card">
                                <h4>Funding Strategy</h4>
                            """, unsafe_allow_html=True)
                            strategy = generate_funding_strategy(startup_idea, industry, "Seed")
                            if strategy:
                                st.write(strategy)
                            st.markdown("</div>", unsafe_allow_html=True)
                else:
                    st.warning("Please describe your startup idea")
        
        # Recent activity section
        st.markdown("""
        <div class="content-card">
            <h4>📈 Your Startup Health Score</h4>
            <p>Track your startup's progress across key metrics:</p>
            <ul>
                <li>Market Position: <strong>Emerging</strong></li>
                <li>Competitive Advantage: <strong>Medium</strong></li>
                <li>Funding Readiness: <strong>High</strong></li>
                <li>Team Strength: <strong>Strong</strong></li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
    
    with tab2:
        st.header("Fundraising Toolkit")
        
        # Investor matching section
        with st.expander("🔍 Smart Investor Matching", expanded=True):
            col1, col2 = st.columns(2)
            with col1:
                funding_stage = st.selectbox("Funding Stage", ["Pre-Seed", "Seed", "Series A", "Series B", "Growth"])
                industry = st.selectbox("Industry", ["Tech", "FinTech", "HealthTech", "EdTech", "AI/ML", "E-commerce"])
            with col2:
                geography = st.selectbox("Geography", ["Global", "North America", "Europe", "Asia", "India"])
                ticket_size = st.selectbox("Ticket Size", ["$100K-$500K", "$500K-$2M", "$2M-$5M", "$5M-$10M", "$10M+"])
            
            if st.button("Find Investors"):
                # Simulate investor matching
                with st.spinner("Matching with ideal investors..."):
                    time.sleep(2)
                    
                    # Display investor cards
                    st.markdown("""
                    <div class="content-card">
                        <h4>Sequoia Capital</h4>
                        <p><strong>Focus:</strong> Early-stage tech, AI/ML, SaaS</p>
                        <p><strong>Recent Investments:</strong> 15 in last 6 months</p>
                        <p><strong>Match Score:</strong> 92%</p>
                        <p><strong>Contact:</strong> partners@sequoiacap.com</p>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    st.markdown("""
                    <div class="content-card">
                        <h4>Accel Partners</h4>
                        <p><strong>Focus:</strong> FinTech, Marketplaces</p>
                        <p><strong>Recent Investments:</strong> 8 in last 6 months</p>
                        <p><strong>Match Score:</strong> 87%</p>
                        <p><strong>Contact:</strong> info@accel.com</p>
                    </div>
                    """, unsafe_allow_html=True)
        
        # Pitch deck analyzer
        with st.expander("📑 AI Pitch Deck Review", expanded=True):
            uploaded_file = st.file_uploader("Upload your pitch deck (PDF or PPTX)", type=["pdf", "pptx"])
            if uploaded_file:
                with st.spinner("Extracting content..."):
                    # Extract text from file
                    text = ""
                    if uploaded_file.type == "application/pdf":
                        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
                        text = "\n".join([page.get_text("text") for page in doc])
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                        prs = Presentation(uploaded_file)
                        text = "\n".join(["\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]) for slide in prs.slides])
                
                if text:
                    review = generate_pitch_deck_review(text)
                    if review:
                        st.markdown(f"""
                        <div class="content-card">
                            <h4>Pitch Deck Analysis</h4>
                            {review}
                        </div>
                        """, unsafe_allow_html=True)
    
    with tab3:
        st.header("Advanced Analytics")
        
        # Market analysis section
        with st.expander("🌍 Market Size Analysis", expanded=True):
            industry = st.selectbox("Select Industry", ["FinTech", "EdTech", "HealthTech", "AI/ML", "E-commerce"])
            region = st.selectbox("Select Region", ["Global", "North America", "Europe", "Asia", "India"])
            
            if st.button("Generate Market Report"):
                with st.spinner("Generating market intelligence..."):
                    time.sleep(3)
                    
                    # Market analysis report
                    st.markdown(f"""
                    <div class="content-card">
                        <h4>Market Analysis: {industry} in {region}</h4>
                        <p><strong>Total Addressable Market:</strong> $12.4B (2024)</p>
                        <p><strong>Growth Rate:</strong> 18.7% CAGR</p>
                        <p><strong>Key Segments:</strong></p>
                        <ul>
                            <li>Segment A: $4.2B (34%)</li>
                            <li>Segment B: $3.1B (25%)</li>
                            <li>Segment C: $2.8B (23%)</li>
                        </ul>
                        <p><strong>Top Players:</strong> Company X (22% share), Company Y (18%), Company Z (12%)</p>
                    </div>
                    """, unsafe_allow_html=True)
        
        # Competitive intelligence
        with st.expander("🕵️ Competitor Intelligence", expanded=True):
            company_name = st.text_input("Enter competitor name")
            if company_name and st.button("Analyze Competitor"):
                with st.spinner("Gathering competitive intelligence..."):
                    time.sleep(3)
                    
                    # Competitor analysis
                    st.markdown(f"""
                    <div class="content-card">
                        <h4>Competitor Analysis: {company_name}</h4>
                        <p><strong>Funding:</strong> Series B ($15M raised)</p>
                        <p><strong>Growth Rate:</strong> 120% YoY</p>
                        <p><strong>Key Metrics:</strong></p>
                        <ul>
                            <li>ARR: $8.2M</li>
                            <li>Customers: 1,250</li>
                            <li>Team Size: 85</li>
                        </ul>
                        <p><strong>Strengths:</strong> Strong brand, efficient CAC</p>
                        <p><strong>Weaknesses:</strong> High churn, concentrated customer base</p>
                    </div>
                    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
